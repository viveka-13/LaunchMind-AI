import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ─── Color Palette ───
INDIGO       = RGBColor(79,  70,  229)   # Primary accent
INDIGO_DARK  = RGBColor(49,  46,  129)   # Darker indigo
BG_LIGHT     = RGBColor(248, 248, 252)   # Slide background
TEXT_DARK    = RGBColor(31,  41,  55)    # Body text
TEXT_MID     = RGBColor(75,  85,  99)    # Secondary text
WHITE        = RGBColor(255, 255, 255)
CYAN         = RGBColor(6,   182, 212)   # Accent highlight

# ─── Slide dimensions (16:9) ───
W = Inches(13.333)
H = Inches(7.5)

# ─── Safe text area constants ───
HEADER_H    = Inches(1.1)
FOOTER_H    = Inches(0.28)
MARGIN      = Inches(0.5)
BODY_TOP    = Inches(1.4)
BODY_H      = H - BODY_TOP - FOOTER_H - Inches(0.15)
BODY_W      = W - MARGIN * 2


def _smart_truncate(text: str, max_chars: int) -> str:
    """Truncate text at a sentence boundary near max_chars."""
    text = str(text).strip()
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars]
    # Try to cut at last newline or period
    last_nl = cut.rfind('\n')
    last_dot = cut.rfind('. ')
    boundary = max(last_nl, last_dot)
    if boundary > max_chars * 0.6:
        return cut[:boundary + 1].strip() + " …"
    return cut.rstrip() + " …"


def _add_bg(slide, prs):
    """Draws background + bottom accent bar."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - FOOTER_H, W, FOOTER_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = INDIGO
    bar.line.fill.background()


def _add_header(slide, title: str, subtitle: str = ""):
    """Draws the indigo header band with title text."""
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, HEADER_H)
    hdr.fill.solid(); hdr.fill.fore_color.rgb = INDIGO
    hdr.line.fill.background()
    # Left accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), HEADER_H)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = CYAN
    stripe.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), W - Inches(0.4), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = title
    run.font.bold = True; run.font.size = Pt(28)
    run.font.color.rgb = WHITE; run.font.name = "Arial"

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.2), Inches(0.72), W - Inches(0.4), Inches(0.3))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run(); run2.text = subtitle
        run2.font.size = Pt(12); run2.font.color.rgb = CYAN; run2.font.name = "Arial"


def _body_text(slide, text: str, left=MARGIN, top=BODY_TOP, width=BODY_W, height=BODY_H,
               font_size=13, color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT):
    """Adds a word-wrapped body text box, hard-limited to the given height."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    # Split into lines for better paragraph control
    lines = text.split('\n')
    first = True
    height_in = int(height) / 914400.0
    width_in = int(width) / 914400.0
    char_budget = int((height_in / 0.18) * (width_in / 0.12))  # rough chars that fit
    used = 0

    for line in lines:
        line = line.strip()
        if not line:
            if not first:
                p = tf.add_paragraph(); p.space_before = Pt(4)
            continue
        if used + len(line) > char_budget:
            break
        used += len(line)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Arial"
        run.font.bold = bold
        p.alignment = align
        p.space_after = Pt(4)
        p.line_spacing = 1.2


def _add_image_safe(slide, img_path: str, left, top, width, height):
    """Adds an image to slide with error handling, maintaining aspect ratio."""
    if not img_path or not os.path.exists(img_path):
        return False
    try:
        from PIL import Image as PILImage
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        aspect = iw / ih
        target_w = int(width)
        target_h = int(int(width) / aspect)
        if target_h > int(height):
            target_h = int(height)
            target_w = int(int(height) * aspect)
        slide.shapes.add_picture(img_path, int(left), int(top), target_w, target_h)
        return True
    except Exception as e:
        print(f"[PPT] Could not add image {img_path}: {e}")
        return False


def generate_pitch_deck_ppt(plan_data: dict, file_path: str,
                             images: list = None, flowchart_path: str = None):
    """
    Generates a professional, well-formatted PowerPoint pitch deck.
    - All text is constrained inside slide boundaries.
    - Pexels images are embedded alongside content.
    - A system architecture flowchart slide is appended at the end.
    """
    if images is None:
        images = []

    os.makedirs(os.path.dirname(file_path) if os.path.dirname(file_path) else '.', exist_ok=True)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]   # Truly blank layout

    img_idx = 0
    def next_img():
        nonlocal img_idx
        if img_idx < len(images):
            p = images[img_idx]; img_idx += 1; return p
        return None

    # ════════════════════════════════════════
    # SLIDE 1 — Title
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    # Full indigo background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = INDIGO_DARK; bg.line.fill.background()
    # Decorative cyan bar top
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = CYAN; top_bar.line.fill.background()
    # Decorative bottom bar
    bot_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.12), W, Inches(0.12))
    bot_bar.fill.solid(); bot_bar.fill.fore_color.rgb = CYAN; bot_bar.line.fill.background()

    # Title text
    title = _smart_truncate(plan_data.get('startup_title', 'Startup Pitch Deck'), 60)
    tb = s.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.333), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run()
    run.text = title; run.font.bold = True
    run.font.size = Pt(54); run.font.color.rgb = WHITE; run.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    run2 = p2.add_run(); run2.text = "Business Plan  |  Generated by AutoStartup AI"
    run2.font.size = Pt(18); run2.font.color.rgb = CYAN; run2.font.name = "Arial"
    p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(16)

    # ════════════════════════════════════════
    # SLIDE 2 — Executive Summary (full-width, no image)
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_bg(s, prs); _add_header(s, "Executive Summary")
    summary = _smart_truncate(plan_data.get('executive_summary', 'N/A'), 600)
    _body_text(s, summary, font_size=14)

    # ════════════════════════════════════════
    # SLIDE 3 — Problem & Target Audience (text left, image right)
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_bg(s, prs)
    _add_header(s, "The Problem & Target Audience")
    img = next_img()
    if img:
        text_w = Inches(7.5)
        content = (f"Problem Statement:\n{_smart_truncate(plan_data.get('problem_statement','N/A'), 280)}"
                   f"\n\nTarget Audience:\n{_smart_truncate(plan_data.get('target_audience','N/A'), 180)}")
        _body_text(s, content, width=text_w, font_size=12)
        _add_image_safe(s, img, Inches(8.1), BODY_TOP, Inches(4.8), BODY_H - Inches(0.2))
    else:
        content = (f"Problem Statement:\n{_smart_truncate(plan_data.get('problem_statement','N/A'), 400)}"
                   f"\n\nTarget Audience:\n{_smart_truncate(plan_data.get('target_audience','N/A'), 280)}")
        _body_text(s, content, font_size=13)

    # ════════════════════════════════════════
    # SLIDE 4 — Solution (image left, text right)
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_bg(s, prs)
    _add_header(s, "The Solution")
    img = next_img()
    if img:
        _add_image_safe(s, img, MARGIN, BODY_TOP, Inches(5.0), BODY_H - Inches(0.2))
        solution = _smart_truncate(plan_data.get('solution', 'N/A'), 400)
        _body_text(s, solution, left=Inches(5.8), width=Inches(7.0), font_size=12)
    else:
        solution = _smart_truncate(plan_data.get('solution', 'N/A'), 700)
        _body_text(s, solution, font_size=13)

    # ════════════════════════════════════════
    # SLIDE 5 — Core Features
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_bg(s, prs)
    _add_header(s, "Core Features")
    img = next_img()
    if img:
        text_w = Inches(7.5)
        features = _smart_truncate(plan_data.get('core_features', 'N/A'), 500)
        _body_text(s, features, width=text_w, font_size=12)
        _add_image_safe(s, img, Inches(8.1), BODY_TOP, Inches(4.8), BODY_H - Inches(0.2))
    else:
        features = _smart_truncate(plan_data.get('core_features', 'N/A'), 700)
        _body_text(s, features, font_size=13)

    # ════════════════════════════════════════
    # SLIDE 6 — Market Gap
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_bg(s, prs)
    _add_header(s, "Market Gap & Opportunity")
    gap = _smart_truncate(plan_data.get('market_gap', 'N/A'), 700)
    _body_text(s, gap, font_size=13)

    # ════════════════════════════════════════
    # SLIDE 7 — Business & Revenue Model
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_bg(s, prs)
    _add_header(s, "Business & Revenue Model")
    revenue = _smart_truncate(plan_data.get('revenue_model', 'N/A'), 700)
    _body_text(s, revenue, font_size=13)

    # ════════════════════════════════════════
    # SLIDE 8 — Implementation Roadmap
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_bg(s, prs)
    _add_header(s, "Implementation Roadmap")
    roadmap = _smart_truncate(plan_data.get('implementation_roadmap', 'N/A'), 700)
    _body_text(s, roadmap, font_size=13)

    # ════════════════════════════════════════
    # SLIDE 9 — System Architecture / Flowchart
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    _add_bg(s, prs)
    _add_header(s, "System Architecture Blueprint",
                subtitle="AI-powered workflow & technical design")

    if flowchart_path and os.path.exists(flowchart_path):
        # Center the flowchart on the slide
        fc_w = Inches(10.0)
        fc_h = Inches(5.5)
        fc_left = int((W - fc_w) / 2)
        fc_top = int(BODY_TOP + Inches(0.1))
        _add_image_safe(s, flowchart_path, fc_left, fc_top, fc_w, fc_h)
    else:
        # Fallback text if flowchart couldn't be generated
        fallback = (
            "System Architecture Overview:\n\n"
            "User Interface → API Gateway → AI Processing Engine\n"
            "AI Engine ↔ Vector Database (ChromaDB) ↔ Knowledge Base\n"
            "AI Engine → Document Generation (PPT / PDF / Word)\n"
            "AI Engine → SQLite (Plan History) → User Dashboard"
        )
        _body_text(s, fallback, font_size=14, color=TEXT_MID)

    # ════════════════════════════════════════
    # SLIDE 10 — Thank You
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = INDIGO; bg.line.fill.background()

    tb = s.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.2))
    tf = tb.text_frame; p = tf.paragraphs[0]; run = p.add_run()
    run.text = "Thank You"; run.font.bold = True
    run.font.size = Pt(54); run.font.color.rgb = WHITE; run.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph(); run2 = p2.add_run()
    run2.text = f"Built with AutoStartup AI  |  {plan_data.get('startup_title', '')}"
    run2.font.size = Pt(18); run2.font.color.rgb = CYAN; run2.font.name = "Arial"
    p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)

    prs.save(file_path)
    print(f"[PPT] Saved: {file_path}  ({len(prs.slides)} slides)")
